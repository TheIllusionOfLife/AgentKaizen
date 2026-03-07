"""Generate AgentKaizen presentation PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_YELLOW = RGBColor(0xFF, 0xBE, 0x0B)
ACCENT_BLUE = RGBColor(0x3A, 0x86, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x06, 0xD6, 0xA0)
RED_ACCENT = RGBColor(0xEF, 0x47, 0x6F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=WHITE, bold=False, space_before=Pt(6),
             alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Meiryo"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)
    return shape


# ============================================================
# Slide 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "AgentKaizen", font_size=54, color=ACCENT_YELLOW, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
             "AIコーディングエージェントの振る舞いを、実験で改善する",
             font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
             "W&B Weave Hackathon",
             font_size=20, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Problem: ステアリング問題", font_size=36, color=ACCENT_YELLOW, bold=True)

# Left side - document list
tf = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5),
                  "コーディングエージェントの振る舞いを制御するドキュメントが多数存在:",
                  font_size=18, color=WHITE)

docs = [
    "AGENTS.md / CLAUDE.md",
    "README.md",
    "Skills / Rules",
    "MCP Servers",
    "Config / Profile",
]
for doc in docs:
    add_para(tf, f"  •  {doc}", font_size=18, color=LIGHT_GRAY, space_before=Pt(8))

# Right side - problem statement
add_rounded_rect(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3.2),
                 RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(2.8),
                  "どのドキュメントを変えれば、出力が良くなるかわからない",
                  font_size=22, color=RED_ACCENT, bold=True)
add_para(tf, "", font_size=10, color=WHITE, space_before=Pt(12))
add_para(tf, "現状は「勘と経験」に頼っている", font_size=18, color=WHITE, space_before=Pt(8))
add_para(tf, "変更の効果を定量的に測定できない", font_size=18, color=WHITE, space_before=Pt(8))
add_para(tf, "ドキュメントが増えるほど特定が困難", font_size=18, color=WHITE, space_before=Pt(8))

# Bottom quote
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1),
             '"AI is easy to demo, hard to productionize"\n'
             '— これはエージェントの設定ドキュメントにも当てはまる',
             font_size=16, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 3: Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Solution: AgentKaizen", font_size=36, color=ACCENT_YELLOW, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "ステアリングドキュメントのためのEvaluation Loop",
             font_size=22, color=WHITE)

# Flow diagram using boxes
steps = [
    ("1. Trace", "エージェント実行を\nWeaveに記録"),
    ("2. Case生成", "実トレースから\nEval Caseを自動生成"),
    ("3. Variant定義", "ドキュメント変更を\nJSONで定義"),
    ("4. Evaluation", "同じDatasetで\nBaseline vs Variant"),
    ("5. Scoring", "Quality / Latency /\nTokenでランキング"),
    ("6. 昇格", "勝ったVariantを\n本番に反映"),
]

box_w = Inches(1.8)
box_h = Inches(2.0)
start_x = Inches(0.5)
y = Inches(2.3)
gap = Inches(0.27)

for i, (title, desc) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    color = ACCENT_BLUE if i < 5 else GREEN
    shape = add_rounded_rect(slide, x, y, box_w, box_h, color)

    add_text_box(slide, x + Inches(0.1), y + Inches(0.2), box_w - Inches(0.2), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.8), box_w - Inches(0.2), Inches(1.0),
                 desc, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Arrow between boxes
    if i < len(steps) - 1:
        arrow_x = x + box_w
        add_text_box(slide, arrow_x, y + Inches(0.7), gap, Inches(0.5),
                     "→", font_size=24, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)

# Bottom tagline
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(1),
             "「勘」から「実験」へ。Vibe Tuningを卒業する。",
             font_size=24, color=ACCENT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.6),
             "One-shot Runner抽象化でCodexとClaude Codeに対応。Eval/セッション分析は現在Codex中心",
             font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 4: Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Demo", font_size=36, color=ACCENT_YELLOW, bold=True)

# Large play button area
add_rounded_rect(slide, Inches(2), Inches(1.5), Inches(9), Inches(4.5),
                 RGBColor(0x2A, 0x2A, 0x4E))

add_text_box(slide, Inches(2), Inches(2.5), Inches(9), Inches(1),
             "▶", font_size=72, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.0), Inches(9), Inches(1),
             "ここで動画を再生します（約2分）",
             font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.6),
             "[YouTubeリンクまたは埋め込み動画]",
             font_size=16, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 5: Weave Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Weave機能の活用マップ", font_size=36, color=ACCENT_YELLOW, bold=True)

# 4 quadrants
quad_w = Inches(5.8)
quad_h = Inches(2.8)
x1, x2 = Inches(0.5), Inches(6.8)
y1, y2 = Inches(1.2), Inches(4.2)

# Trace
add_rounded_rect(slide, x1, y1, quad_w, quad_h, RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, x1 + Inches(0.3), y1 + Inches(0.15), quad_w - Inches(0.6), quad_h - Inches(0.3),
                  "Trace (トレーシング)", font_size=20, color=ACCENT_BLUE, bold=True)
add_para(tf, "• One-shot実行のInput/Output/メタデータ自動記録", font_size=13, color=WHITE, space_before=Pt(10))
add_para(tf, "• インタラクティブセッション取り込み・再構築", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "• マルチモーダル: 画像プロンプトのcontent block保持", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "• PIIリダクション: Weave組み込み + カスタム", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "  (APIキー、ローカルパス、ユーザー名の除去)", font_size=12, color=LIGHT_GRAY, space_before=Pt(2))

# Asset Management
add_rounded_rect(slide, x2, y1, quad_w, quad_h, RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, x2 + Inches(0.3), y1 + Inches(0.15), quad_w - Inches(0.6), quad_h - Inches(0.3),
                  "Asset Management (アセット管理)", font_size=20, color=GREEN, bold=True)
add_para(tf, "• Dataset: JSONL Eval Case + トレースからの自動生成", font_size=13, color=WHITE, space_before=Pt(10))
add_para(tf, "• Scorer: カスタム10種 + Weave組み込み", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "  (ValidJSON, Pydantic)", font_size=12, color=LIGHT_GRAY, space_before=Pt(2))
add_para(tf, "• Model: CodexVariantModel (weave.Model)", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "  ワークスペース構成+設定の組み合わせをバージョン管理", font_size=12, color=LIGHT_GRAY, space_before=Pt(2))

# Offline Evaluation
add_rounded_rect(slide, x1, y2, quad_w, quad_h, RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, x1 + Inches(0.3), y2 + Inches(0.15), quad_w - Inches(0.6), quad_h - Inches(0.3),
                  "Offline Evaluation (オフライン評価)", font_size=20, color=ACCENT_BLUE, bold=True)
add_para(tf, "• weave.Evaluationで同じDatasetに対しVariant比較", font_size=13, color=WHITE, space_before=Pt(10))
add_para(tf, "• 各Scorerの結果を横並びで確認・ランキング", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "• コスト・レイテンシ追跡:", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "  品質同等でもコスト増ならGate不合格", font_size=12, color=LIGHT_GRAY, space_before=Pt(2))

# Online Evaluation
add_rounded_rect(slide, x2, y2, quad_w, quad_h, RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, x2 + Inches(0.3), y2 + Inches(0.15), quad_w - Inches(0.6), quad_h - Inches(0.3),
                  "Online Evaluation的活用", font_size=20, color=GREEN, bold=True)
add_para(tf, "• One-shot実行時にScorerをリアルタイム適用", font_size=13, color=WHITE, space_before=Pt(10))
add_para(tf, "• セッションのヒューリスティクスコアリング:", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "  ブランチ作成・テスト実行・lint実行を自動検出", font_size=12, color=LIGHT_GRAY, space_before=Pt(2))
add_para(tf, "• 今後: Weave Monitor/Guardrailとの統合", font_size=13, color=MEDIUM_GRAY, space_before=Pt(4))

# ============================================================
# Slide 6: Custom Implementation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "カスタム実装の工夫", font_size=36, color=ACCENT_YELLOW, bold=True)

# Left column
col_w = Inches(5.8)
left_x = Inches(0.5)
right_x = Inches(6.8)

# 1. Scorers
add_rounded_rect(slide, left_x, Inches(1.2), col_w, Inches(2.5), RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, left_x + Inches(0.3), Inches(1.3), col_w - Inches(0.6), Inches(2.3),
                  "Deterministic Scorer群 (10種)", font_size=18, color=ACCENT_BLUE, bold=True)
add_para(tf, "• 必須テキスト含有 / 禁止テキスト不在 / 完全一致", font_size=13, color=WHITE, space_before=Pt(8))
add_para(tf, "• 文字数制限 (最小・最大) / JSON妥当性", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "• セクション存在 / コンテンツグループ網羅", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "• ファイルパス引用 / トークン使用量", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "→ LLM不要の高速・安定・再現性のあるチェック", font_size=13, color=GREEN, space_before=Pt(8))

# 2. Variant Workspace
add_rounded_rect(slide, right_x, Inches(1.2), col_w, Inches(2.5), RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, right_x + Inches(0.3), Inches(1.3), col_w - Inches(0.6), Inches(2.3),
                  "Variant Workspace管理", font_size=18, color=ACCENT_BLUE, bold=True)
add_para(tf, "• リポジトリを一時コピー → 変更適用 → 同条件で実行", font_size=13, color=WHITE, space_before=Pt(8))
add_para(tf, "• JSON定義: append / prepend / replace", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "• 外部ファイル注入にも対応", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "→ 本番リポジトリを汚さず安全に実験", font_size=13, color=GREEN, space_before=Pt(8))

# 3. Ranking & Gate
add_rounded_rect(slide, left_x, Inches(4.0), col_w, Inches(2.5), RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, left_x + Inches(0.3), Inches(4.1), col_w - Inches(0.6), Inches(2.3),
                  "Quality Ranking & Regression Gate", font_size=18, color=ACCENT_BLUE, bold=True)
add_para(tf, "• Quality Score = アクティブScorerの加重Pass率", font_size=13, color=WHITE, space_before=Pt(8))
add_para(tf, "• Gate: 品質同等でもLatency/Token増 → 不合格", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "→ 「良くなったけど遅くなった」変更を自動検出", font_size=13, color=GREEN, space_before=Pt(8))

# 4. Session Analysis & Case Gen
add_rounded_rect(slide, right_x, Inches(4.0), col_w, Inches(2.5), RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, right_x + Inches(0.3), Inches(4.1), col_w - Inches(0.6), Inches(2.3),
                  "セッション分析 & Case自動生成", font_size=18, color=ACCENT_BLUE, bold=True)
add_para(tf, "• ヒューリスティクス + 外部Judge の2層スコアリング", font_size=13, color=WHITE, space_before=Pt(8))
add_para(tf, "• optimization_relevance: 次に直すべき", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "  ドキュメントを推薦 (AGENTS/README/Skill/Config)", font_size=12, color=LIGHT_GRAY, space_before=Pt(2))
add_para(tf, "• Weave Call History APIから過去トレース →", font_size=13, color=WHITE, space_before=Pt(4))
add_para(tf, "  重複排除 → JSONL Eval Case自動生成", font_size=12, color=LIGHT_GRAY, space_before=Pt(2))

# ============================================================
# Slide 7: Example Experiment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "実験例: AGENTS.mdで日本語レスポンス制御", font_size=32, color=ACCENT_YELLOW, bold=True)

# Variant definition
add_rounded_rect(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.2), RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.3), Inches(2.0),
                  "Variant定義 (JSON)", font_size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, '{ "name": "agents-japanese-response",', font_size=13, color=LIGHT_GRAY, space_before=Pt(8),
         font_name="Courier New")
add_para(tf, '  "edits": [{ "path": "AGENTS.md",', font_size=13, color=LIGHT_GRAY, space_before=Pt(2),
         font_name="Courier New")
add_para(tf, '    "mode": "append",', font_size=13, color=LIGHT_GRAY, space_before=Pt(2),
         font_name="Courier New")
add_para(tf, '    "text": "You must respond in Japanese."', font_size=13, color=ACCENT_YELLOW, space_before=Pt(2),
         font_name="Courier New")
add_para(tf, '  }] }', font_size=13, color=LIGHT_GRAY, space_before=Pt(2),
         font_name="Courier New")

# Results table area
add_rounded_rect(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.2), RGBColor(0x2A, 0x2A, 0x4E))
tf = add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.3), Inches(2.0),
                  "結果", font_size=16, color=GREEN, bold=True)
add_para(tf, "✓ 日本語プロンプト → 日本語で回答", font_size=14, color=WHITE, space_before=Pt(10))
add_para(tf, '  must_contain: "リポジトリ", "W&B Weave"', font_size=12, color=LIGHT_GRAY, space_before=Pt(2))
add_para(tf, '✓ "Say only: ok" → "ok" (過剰適用なし)', font_size=14, color=WHITE, space_before=Pt(8))
add_para(tf, '  exact_match で検証', font_size=12, color=LIGHT_GRAY, space_before=Pt(2))
add_para(tf, '✓ "Respond in English" → 英語 (ユーザー指示優先)', font_size=14, color=WHITE, space_before=Pt(8))

# Bottom summary
add_rounded_rect(slide, Inches(0.5), Inches(3.8), Inches(12.1), Inches(1.6), RGBColor(0x1E, 0x3A, 0x5F))
tf = add_text_box(slide, Inches(0.8), Inches(3.9), Inches(11.5), Inches(1.4),
                  "Evaluation結果", font_size=18, color=ACCENT_YELLOW, bold=True)
add_para(tf, "• Quality Score向上: Variantが日本語ケースでBaselineを上回る", font_size=15, color=WHITE, space_before=Pt(8))
add_para(tf, "• Gate Pass: Latency/Token数のリグレッションなし → 昇格可能", font_size=15, color=WHITE, space_before=Pt(4))
add_para(tf, "• コントロールケースで副作用がないことも定量的に確認", font_size=15, color=WHITE, space_before=Pt(4))

# ============================================================
# Slide 8: Future
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "今後の展望", font_size=36, color=ACCENT_YELLOW, bold=True)

items = [
    ("PyPIリリース", "pip install agentkaizen で誰でも利用可能に", "進行中"),
    ("Claude Code完全対応", "セッション分析、トークン使用量取得", "進行中"),
    ("Semantic Scorer", "LLM-as-a-Judgeで回答の意味的品質も評価", "計画中"),
    ("マルチモーダルEval拡充", "画像入力ケースの評価スコアリング", "計画中"),
    ("Online Evaluation連携", "Weave Monitor/Guardrailとの統合", "計画中"),
    ("CI/CD統合", "PRごとにドキュメント変更のEvalを自動実行", "計画中"),
]

for i, (title, desc, status) in enumerate(items):
    y = Inches(1.3) + i * Inches(0.85)
    status_color = GREEN if status == "進行中" else ACCENT_BLUE
    add_rounded_rect(slide, Inches(0.5), y, Inches(12.1), Inches(0.7), RGBColor(0x2A, 0x2A, 0x4E))

    add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(3.5), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(4.5), y + Inches(0.08), Inches(6), Inches(0.5),
                 desc, font_size=15, color=LIGHT_GRAY)
    add_text_box(slide, Inches(10.8), y + Inches(0.08), Inches(1.5), Inches(0.5),
                 status, font_size=14, color=status_color, bold=True, alignment=PP_ALIGN.RIGHT)

# Tagline
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8),
             "AgentKaizen: Measure changes, don't guess.\n「勘」ではなく「実験」でエージェントを改善する。",
             font_size=22, color=ACCENT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

# Save
output_path = "AgentKaizen_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
